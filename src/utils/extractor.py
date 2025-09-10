import os
import base64
import subprocess
from pathlib import Path
from io import BytesIO
from PIL import Image
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
import tempfile


def encode_image_to_base64(image: Image.Image, format="PNG") -> str:
    """Encode PIL image to base64 string"""
    buffer = BytesIO()
    image.save(buffer, format=format)
    return base64.b64encode(buffer.getvalue()).decode("utf-8")


def libreoffice_to_pdf(input_path: str, output_dir: str) -> str:
    """Convert any supported document to PDF using LibreOffice (headless)."""
    cmd = [
        "soffice",
        "--headless",
        "--convert-to", "pdf",
        "--outdir", output_dir,
        input_path,
    ]
    subprocess.run(cmd, check=True)
    pdf_file = os.path.splitext(os.path.basename(input_path))[0] + ".pdf"
    return os.path.join(output_dir, pdf_file)


def process_document(file_path,
                     dpi=200,
                     image_format="PNG",
                     processing_mode="quality"):
    """
    Process PDF, DOCX, DOC, TXT, PPT, PPTX and image files.

    processing_mode:
        - "normal": extracts text whenever possible
        - "quality": converts everything to images via LibreOffice + PyMuPDF

    Returns a list of items in the form:
      {"type": "image", "content": <base64>}
      {"type": "text", "content": <string>}
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")

    ext = Path(file_path).suffix.lower()
    results = []

    # ---------------- QUALITY MODE ----------------
    if processing_mode == "quality":
        if ext == ".pdf":
            pdf_path = file_path
        else:
            with tempfile.TemporaryDirectory() as tmpdir:
                pdf_path = libreoffice_to_pdf(file_path, tmpdir)
                doc = fitz.open(pdf_path)
                for page in doc:
                    pix = page.get_pixmap(dpi=dpi)
                    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    b64 = encode_image_to_base64(img, format=image_format.upper())
                    results.append({"type": "image", "content": b64})
            return results

        # Caso seja PDF, processa direto (sem criar tmpdir)
        doc = fitz.open(pdf_path)
        for page in doc:
            pix = page.get_pixmap(dpi=dpi)
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            b64 = encode_image_to_base64(img, format=image_format.upper())
            results.append({"type": "image", "content": b64})
        return results

    # ---------------- NORMAL MODE ----------------
    if ext == ".pdf":
        doc = fitz.open(file_path)
        for page in doc:
            text = page.get_text().strip()
            if text:
                for line in text.splitlines():
                    if line.strip():
                        results.append({"type": "text", "content": line.strip()})

    elif ext == ".docx":
        doc = Document(file_path)
        for para in doc.paragraphs:
            if para.text.strip():
                results.append({"type": "text", "content": para.text.strip()})

    elif ext == ".doc":
        try:
            result = subprocess.run(
                ["antiword", file_path],
                capture_output=True,
                text=True,
                check=True
            )
            text = result.stdout.strip()
            if text:
                for line in text.splitlines():
                    if line.strip():
                        results.append({"type": "text", "content": line.strip()})
        except FileNotFoundError:
            raise RuntimeError("Antiword is not installed or not in PATH.")
        except subprocess.CalledProcessError as e:
            raise RuntimeError(f"Antiword failed to extract text: {e}")

    elif ext == ".txt":
        with open(file_path, "r", encoding="utf-8") as f:
            for line in f:
                if line.strip():
                    results.append({"type": "text", "content": line.strip()})

    elif ext in [".ppt", ".pptx"]:
        prs = Presentation(file_path)
        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            if slide_text:
                results.append({"type": "text", "content": "\n".join(slide_text)})

    elif ext in [".png", ".jpg", ".jpeg", ".bmp", ".gif", ".webp", ".tiff", ".svg"]:
        img = Image.open(file_path)
        b64 = encode_image_to_base64(img, format=image_format.upper())
        results.append({"type": "image", "content": b64})

    else:
        raise ValueError(f"Unsupported file format: {ext}")

    return results